import collections 
import collections.abc
from pptx import Presentation 
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
# import pptx.dml.color.ColorFormat
import json
from pptx.util import Inches, Pt
import os
from pptx.enum.shapes import MSO_SHAPE
import requests
import base64
import shutil
import copy
import random
import six
from pptx.dml.color import RGBColor
import hashlib
import pptx
import time
import pathlib
from replacement import change_image_of_slide, generate_initial_prs, replace_text, change_image_and_text_on_slide
from pptx.enum.text import PP_ALIGN
from googletrans import Translator
import gpt4free.g4f as g4f
from generateImages import generateImage

g4f.debug.logging = True
g4f.check_version = False

translator = Translator()

num_of_imgs = ['1', '1', '2']
path_to_slides_parser = pathlib.Path().resolve()
path_to_slides_parser = os.path.abspath(os.path.join(path_to_slides_parser, os.pardir))
print(path_to_slides_parser)

def question(history=[{}]):
    answer = ""
    response = g4f.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=history,
        stream=True,
    )

    for message in response:
        answer += message
    return answer


def separated_generateThemes(numOfSlides, topic):
    try:
        final_answer = []
        initial_answer = question( 
                                history=[{"role": "user", "content" : "Im creating a presentation about calculus, generate 3 themes for 3 slides, then generate prompt for picture to this slide, then generate text with 40 words about this lide. Devide every part by '-' without any additional words"}, 
                                        {"role": "assistant", "content" : 'Introduction to Calculus - A graph with a curved line representing a function - "Calculus is a branch of mathematics that deals with change and motion. It provides tools to analyze and understand how quantities change over time. By studying rates of change and the accumulation of quantities." - Differentiation - A tangent line to a curve at a specific point - "Differentiation is a fundamental concept in calculus. By calculating derivatives, we can determine the slope of a curve, identify maximum and minimum points, and analyze the behavior of functions. Differentiation plays a crucial role in optimization and modeling real-world phenomena." - Integration - The area under a curve bounded by the x-axis - "Integration is another key concept in calculus. It focuses on finding the total accumulation of a quantity over a given interval. By calculating definite or indefinite integrals, we can determine the area under a curve, solve problems related to displacement, velocity, and acceleration, and evaluate complex mathematical expressions."'},
                        {"role": "user", "content": f"Im creating a presentation about {topic}, generate {numOfSlides-1} themes for {numOfSlides-1} slides, then generate prompt for picture to this slide, then generate text with 40 words about this lide. Devide every part by '-' without any additional words"}])
        splitted_answer = initial_answer.replace('\n', ' - ').replace('"', '').split(' - ')
        for i in range(0, len(splitted_answer), 3):
            final_answer.append([splitted_answer[i], 'sdxl', splitted_answer[i+1], random.choice(num_of_imgs), splitted_answer[i+1], splitted_answer[i+2]])
        return final_answer
    except:
        return [[]]


def firstSlide(slide, topic, usersname):
    top = Pt(400)
    width = Pt(1920)
    left = Pt(0)
    height = Pt(150)
    txBox = slide.shapes.add_textbox(left, top, width, height)# (left, top, width, height)

    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Times New Roman' #Arial Rounded MT Bold
    p.font.size = Pt(160)
    p = tf.add_paragraph()
    p.text = "Prepared by " + usersname
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.name = 'Times New Roman'
    p.font.color.rgb = RGBColor(0, 0, 0)
    print("1st Slide done")



# themes = [['Anatomy', 'dalle', 'Frog anatomy diagram', '2', "A diagram of a frog's anatomy, highlighting its unique features such as webbed feet and long tongue", 'Frogs have unique adaptations that allow them to thrive in their aquatic and terrestrial habitats. This diagram showcases the anatomy of a frog, including its webbed feet, long tongue, and bulging eyes. '], ['Habitat', 'network', 'Frog habitat image', '1', 'An image of a frog in its natural habitat, such as a pond or wetland', 'Frogs are found in a variety of habitats, from rainforests to deserts. This image captures a frog in its natural habitat, showcasing the importance of wetlands and other aquatic ecosystems for frog survival. '], ['Life cycle', 'dalle', 'Frog life cycle illustration', '2', "An illustration of the different stages of a frog's life cycle, from egg to tadpole to adult frog", "Frogs undergo a unique metamorphosis, transforming from aquatic tadpoles to terrestrial adults. This illustration depicts the different stages of a frog's life cycle, highlighting the importance of wetland habitats for breeding and development. "], ['Diet', 'network', 'Frog eating insect', '1', 'An image of a frog catching and eating an insect, showcasing its role as a predator in the food chain', 'Frogs are important predators in their ecosystems, feeding on insects and other small animals. This image captures a frog in action, highlighting its unique feeding behavior and role in maintaining a healthy ecosystem. '], ['Conservation', 'dalle', 'Endangered frog species', '3', 'A collection of images showcasing endangered frog species and the threats they face, such as habitat loss and pollution', 'Frogs are facing numerous threats, including habitat loss, pollution, and disease. This collection of images highlights some of the endangered frog species and the urgent need for conservation efforts to protect these important amphibians.']]



def generate_presentation(main_topic, usersname, num_of_slides, filename, timecode=1, color_of_theme="black"):
    one_picture_slides_for_dark = [1, 2, 3, 4, 5, 6, 7, 8, 9]
    two_picture_slides_for_dark = [10, 11, 12, 13, 14]

    one_picture_slides_for_white = [16, 17, 18, 19, 20, 21, 22, 23, 24]
    two_picture_slides_for_white = [25, 26, 27, 28, 29]
    
    language = translator.detect(main_topic).lang
# 
    # language = translator.detect(str(main_topic)).lang
    
    print("start presentation")
    for i in range(5):
        counter_of_mistakes_in_themes = 0
        themes = separated_generateThemes(num_of_slides, translator.translate(main_topic, dest='en').text)
        # themes = [['Introduction to Tenge', 'sdxl', 'An image of the Kazakhstani tenge banknotes and coins', '1', 'An image of the Kazakhstani tenge banknotes and coins', "Tenge is the national currency of Kazakhstan, introduced in 1993 after the country gained independence. It plays a crucial role in the country's economy and is widely used for daily transactions. The banknotes feature various important landmarks and symbols of Kazakhstan, while the coins represent different denominations."], ['History of Tenge', 'sdxl', 'A timeline showcasing the evolution of tenge banknotes', '1', 'A timeline showcasing the evolution of tenge banknotes', "The history of tenge dates back to ancient times when trade and bartering systems were prevalent. Over the years, the currency has undergone several changes, including design modifications, currency reforms, and stability measures. Today, the tenge embodies the nation's rich history and serves as a symbol of economic growth and development."], ['Tenge Exchange Rates', 'sdxl', 'A comparison chart displaying exchange rates with other currencies', '2', 'A comparison chart displaying exchange rates with other currencies', 'Understanding the exchange rates of tenge is essential for international trade and travel. The value of tenge fluctuates based on various factors such as supply and demand, inflation rates, and global economic conditions. Keeping track of exchange rates helps businesses, tourists, and investors make informed decisions regarding their financial transactions involving the Kazakhstani currency.'], ['Security Features of Tenge', 'sdxl', 'An image highlighting the security features on tenge banknotes', '1', 'An image highlighting the security features on tenge banknotes', 'Tenge banknotes incorporate advanced security features to prevent counterfeiting and ensure the integrity of the currency. These features include holograms, watermarks, microprinting, and special inks. By incorporating these elements, the National Bank of Kazakhstan aims to maintain trust in the currency and protect it from fraudulent activities.'], ['Future Outlook of Tenge', 'sdxl', 'An illustration depicting potential advancements in the use of tenge', '2', 'An illustration depicting potential advancements in the use of tenge', 'As technology continues to evolve, the future of tenge holds exciting possibilities. With the emergence of digital currencies, there is a growing interest in exploring blockchain-based solutions for payments using tenge. The National Bank of Kazakhstan is actively researching and adopting innovative technologies to enhance the efficiency, security, and usability of the currency in the digital age.']]
        # themes = [['Introduction to Tenge', 'sdxl', 'An image of the Kazakhstani tenge currency', '1', 'An image of the Kazakhstani tenge currency', "The tenge is the official currency of Kazakhstan. It was introduced in 1993 after the dissolution of the Soviet Union. The tenge is represented by the symbol ₸ and is subdivided into 100 tiyn. It plays a crucial role in the country's economy and is widely used for both domestic and international transactions. "], ['History of Tenge', 'sdxl', 'A timeline showcasing the evolution of the tenge', '1', 'A timeline showcasing the evolution of the tenge', "The history of the tenge dates back to the 15th century when various forms of currency were used in the region. However, it was not until 1993 that the modern tenge was introduced as the official currency of Kazakhstan. Over the years, the tenge has undergone several changes in design and denominations, reflecting the country's economic development. "], ['Security features of Tenge', 'sdxl', 'A close-up image highlighting the security features of the tenge banknotes', '1', 'A close-up image highlighting the security features of the tenge banknotes', 'The tenge banknotes incorporate various security features to prevent counterfeiting and ensure the integrity of the currency. These features include watermarks, holograms, microprinting, and special inks. The National Bank of Kazakhstan continuously updates and enhances these security measures to maintain the trust and confidence of the public in the tenge. '], ['Exchange rates and international use of Tenge', 'sdxl', 'A world map showing countries where the tenge is accepted or commonly exchanged', '1', 'A world map showing countries where the tenge is accepted or commonly exchanged', 'The tenge is primarily used within Kazakhstan, but it can also be exchanged in several neighboring countries and international financial centers. The exchange rate of the tenge fluctuates based on various factors such as economic conditions, inflation, and global market trends. Understanding the exchange rates is essential for businesses and individuals involved in international trade and travel. '], ['Future prospects of Tenge', 'sdxl', 'An image depicting the potential growth and development of the Kazakhstani economy', '1', 'An image depicting the potential growth and development of the Kazakhstani economy', 'The future prospects of the tenge are closely tied to the economic growth and stability of Kazakhstan. As the country continues to diversify its economy and attract foreign investments, the value and importance of the tenge are expected to strengthen. The government and the National Bank of Kazakhstan play a crucial role in ensuring the stability and sustainability of the tenge in the global financial landscape.']]
        print(themes)
        print("FIRST ENDPOINT")
        if len(themes) < num_of_slides-1:
            counter_of_mistakes_in_themes = 1
        print("Second ENDPOINT")

        for j in themes:
            for k in j:
                if k == "":
                    counter_of_mistakes_in_themes += 1
        print("Third ENDPOINT")

        if counter_of_mistakes_in_themes == 0:
            break
    print("themes approved")
    # Initial
    prs = Presentation()
    prs.slide_height = Pt(1080)
    prs.slide_width = Pt(1920)
    slides = []
    which_slides = [15]
    if color_of_theme=="black":
        which_slides = [0]
    
    
    blank_slide_layout = prs.slide_layouts[6]
    main_root = f"{path_to_slides_parser}/imgs/"
    direction = f"{main_root}{str(usersname.replace(' ', ''))}{str(timecode)}"
    try:
        os.mkdir(direction)
    except:
        pass
    print("Directory created")

    for i in range(len(themes)):
        try:
            if color_of_theme=="black":
                if str(themes[i][3]) == "2":
                    slide_number_from_list = random.choice(two_picture_slides_for_dark)
                    which_slides.append(slide_number_from_list)
                    # two_picture_slides_for_dark.remove(slide_number_from_list)
                else:
                    slide_number_from_list = random.choice(one_picture_slides_for_dark)
                    which_slides.append(slide_number_from_list)
                    one_picture_slides_for_dark.remove(slide_number_from_list)
            else:
                if str(themes[i][3]) == "2":
                    slide_number_from_list = random.choice(two_picture_slides_for_white)
                    which_slides.append(slide_number_from_list)
                    # two_picture_slides_for_white.remove(slide_number_from_list)
                else:
                    slide_number_from_list = random.choice(one_picture_slides_for_white)
                    which_slides.append(slide_number_from_list)
                    one_picture_slides_for_white.remove(slide_number_from_list)
        except Exception as e:
            print(e)
            slide_number_from_list = random.choice(one_picture_slides_for_dark)
            which_slides.append(slide_number_from_list)
            one_picture_slides_for_dark.remove(slide_number_from_list)
    print("slides decided")

    prs = generate_initial_prs(Presentation(rf"{path_to_slides_parser}/slidesTemplate/Final_Template.pptx"), which_slides, filename) #FIX (filename in function donr require)
    replace_text({"Theme":main_topic, "username":f'Prepared by {usersname}'}, prs.slides[0])
    flag = True
    
    print("initial presentation created")

    for i in range(1, len(themes)+1):
        print(f"{i}'s slide done")
        images_paths = []
        # Create an images
        for j in range(int(themes[i-1][3])):
            # generateImage(name=f"image{str(i)+str(j)}", dir=f"{str(usersname.replace(' ', ''))}{str(timecode)}", prompt=themes[j-1][2]) #FIX prompt=themes[j-1][2]
            generateImage(name=f"image{str(i)+str(j)}", dir=f"{direction}", prompt=themes[i-1][4])
            # with open(rf"{direction}\image{str(i)+str(j)}.jpg", 'wb') as handler:
            #     handler.write(img_data)
            if flag:
                shutil.copyfile(str(f"{direction}/image{str(i)+str(j)}.jpeg"), str(f"{path_to_slides_parser}/backend/static/images/const/{str(usersname)+str(timecode)}.jpeg"))
                flag = False
            images_paths.append(rf"{direction}\image{str(i)+str(j)}.jpeg")
        # Change in slide
        print(images_paths)
        change_image_and_text_on_slide(prs, prs.slides[int(i)], images_paths, {"Topic":str(translator.translate(themes[i-1][0], dest=language).text), "Text":str(translator.translate(themes[i-1][5], dest=language).text)})
        # Delete images
        print("text and images changed")
        for j in images_paths:
            # print(1)
            os.remove(j)
    prs.save(f'./static/presentations/{filename}.pptx')
    os.remove(rf'{filename}.pptx')
    # os.remove(rf'./slides_parser/backend/{filename}.pptx')
    os.rmdir(direction)
    try:
        os.rmdir(f'{path_to_slides_parser}/imgs/{str(usersname)+str(timecode)}')
    except:
        pass

# generate_presentation("Frogs", "Bogdan", 2,"Lastpyqq", time.time())