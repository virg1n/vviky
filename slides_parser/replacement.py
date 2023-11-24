import collections 
import collections.abc
from pptx import Presentation 
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
# import pptx.dml.color.ColorFormat
import json
from pptx.util import Inches, Pt
import os
from pptx.enum.shapes import MSO_SHAPE
import requests
import base64
import shutil
import copy
import random
import six
from pptx.dml.color import RGBColor
import hashlib
import pptx
from pptx.opc.package import _Relationship
from pptx.opc.constants import RELATIONSHIP_TARGET_MODE as RTM

themes = [['Anatomy', 'dalle', 'Frog anatomy diagram', '2', "A diagram of a frog's anatomy, highlighting its unique features such as webbed feet and long tongue", 'Frogs have unique adaptations that allow them to thrive in their aquatic and terrestrial habitats. This diagram showcases the anatomy of a frog, including its webbed feet, long tongue, and bulging eyes. '], ['Habitat', 'network', 'Frog habitat image', '1', 'An image of a frog in its natural habitat, such as a pond or wetland', 'Frogs are found in a variety of habitats, from rainforests to deserts. This image captures a frog in its natural habitat, showcasing the importance of wetlands and other aquatic ecosystems for frog survival. '], ['Life cycle', 'dalle', 'Frog life cycle illustration', '2', "An illustration of the different stages of a frog's life cycle, from egg to tadpole to adult frog", "Frogs undergo a unique metamorphosis, transforming from aquatic tadpoles to terrestrial adults. This illustration depicts the different stages of a frog's life cycle, highlighting the importance of wetland habitats for breeding and development. "], ['Diet', 'network', 'Frog eating insect', '1', 'An image of a frog catching and eating an insect, showcasing its role as a predator in the food chain', 'Frogs are important predators in their ecosystems, feeding on insects and other small animals. This image captures a frog in action, highlighting its unique feeding behavior and role in maintaining a healthy ecosystem. '], ['Conservation', 'dalle', 'Endangered frog species', '3', 'A collection of images showcasing endangered frog species and the threats they face, such as habitat loss and pollution', 'Frogs are facing numerous threats, including habitat loss, pollution, and disease. This collection of images highlights some of the endangered frog species and the urgent need for conservation efforts to protect these important amphibians.']]

def _exp_add_slide(ppt, slide_layout):
    """
    Function to handle slide creation in the Presentation, to avoid issues caused by default implementation.

    :param slide_layout:
    :return:
    """

    def generate_slide_partname(self):
        """Return |PackURI| instance containing next available slide partname."""
        from pptx.opc.packuri import PackURI

        sldIdLst = self._element.get_or_add_sldIdLst()

        existing_rels = [f.target_partname for f in self.rels]
        partname_str = "/ppt/slides/slide%d.xml" % (len(sldIdLst) + 1)

        while partname_str in existing_rels:
            import random
            import string

            random_part = ''.join(random.choice(string.ascii_letters) for i in range(2))
            partname_str = "/ppt/slides/slide%s%d.xml" % (random_part, len(sldIdLst) + 1)

        return PackURI(partname_str)

    def add_slide_part(self, slide_layout):
        """
        Return an (rId, slide) pair of a newly created blank slide that
        inherits appearance from *slide_layout*.
        """
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.parts.slide import SlidePart

        partname = generate_slide_partname(self)
        slide_layout_part = slide_layout.part
        slide_part = SlidePart.new(partname, self.package, slide_layout_part)
        rId = self.relate_to(slide_part, RT.SLIDE)
        return rId, slide_part.slide

    def add_slide_ppt(self, slide_layout):
        rId, slide = add_slide_part(self.part, slide_layout)
        slide.shapes.clone_layout_placeholders(slide_layout)
        self._sldIdLst.add_sldId(rId)
        return slide

    # slide_layout = self.get_master_slide_layout(slide_layout)
    return add_slide_ppt(ppt.slides, slide_layout)


def copy_shapes(source, dest):
    """
    Helper to copy shapes handling edge cases.

    :param source:
    :param dest:
    :return:
    """
    from pptx.shapes.group import GroupShape
    import copy

    # Copy all existing shapes
    for shape in source:
        if isinstance(shape, GroupShape):
            group = dest.shapes.add_group_shape()
            group.name = shape.name
            group.left = shape.left
            group.top = shape.top
            group.width = shape.width
            group.height = shape.height
            group.rotation = shape.rotation

            # Recursive copy of contents
            copy_shapes(shape.shapes, group)

            # Fix offset
            cur_el = group._element.xpath(".//p:grpSpPr")[0]
            ref_el = shape._element.xpath(".//p:grpSpPr")[0]
            parent = cur_el.getparent()
            parent.insert(
                parent.index(cur_el) + 1,
                copy.deepcopy(ref_el)
            )
            parent.remove(cur_el)

            result = group
        elif hasattr(shape, "image"):
            import io

            # Get image contents
            content = io.BytesIO(shape.image.blob)
            result = dest.shapes.add_picture(
                content, shape.left, shape.top, shape.width, shape.height
            )
            result.name = shape.name
            result.crop_left = shape.crop_left
            result.crop_right = shape.crop_right
            result.crop_top = shape.crop_top
            result.crop_bottom = shape.crop_bottom
        elif hasattr(shape, "has_chart") and shape.has_chart:
            from .charts import clone_chart
            result = clone_chart(shape, dest)
        else:
            import copy

            newel = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(newel, "p:extLst")
            result = dest.shapes[-1]


# def duplicate_slide(ppt, slide_index: int):
    """
    Duplicate the slide with the given number in presentation.
    Adds the new slide by default at the end of the presentation.

    :param ppt:
    :param slide_index: Slide number
    :return:
    """
    source = ppt.slides[slide_index]

    dest = _exp_add_slide(ppt, source.slide_layout)
    # ppt.slides.add_slide(dest)

    # Remove all shapes from the default layout
    for shape in dest.shapes:
        remove_shape(shape)

    # Copy all existing shapes
    copy_shapes(source.shapes, dest)

    # Copy all existing shapes
    if source.has_notes_slide:
        txt = source.notes_slide.notes_text_frame.text
        dest.notes_slide.notes_text_frame.text = txt

    return dest

def remove_shape(shape):

    el = shape.element  # --- get reference to XML element for shape
    el.getparent().remove(el)  # --- remove that shape element from its tree

def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]

def duplicate_slide(pres, index):
    """Duplicate the slide with the given index in pres.

    Adds slide to the end of the presentation"""
    source = pres.slides[index]

    blank_slide_layout = _get_blank_slide_layout(pres)
    dest = pres.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for key, value in six.iteritems(source.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if not "notesSlide" in value.reltype:
            # dest.part.rels._add_relationship(value.reltype, value._target, value.rId)
            dest.part.rels._rels[value.rId] = _Relationship(
                dest.part.rels._base_uri,
                value.rId,
                value.reltype,
                target_mode=RTM.EXTERNAL if value.is_external else RTM.INTERNAL,
                target=value._target,
            )

    return dest


def generate_initial_prs(prs, slides, final_name):
    for i in slides:
        duplicate_slide(prs, i)
        # pass

    for i in range(len(prs.slides)-1-len(slides), -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    prs.save(f"{final_name}.pptx")
    return prs


def change_image_of_slide(new_images, slide):
    k = 0
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                smallfile = new_images[k]   #r"C:\Users\Bogdan\java\pyth\t2p\presentations\server\imgs\test1\Dog1.png"
                new_pptx_img = pptx.parts.image.Image.from_file(smallfile)
                
                slide_part, rId = shape.part, shape._element.blip_rId
                image_part = slide_part.related_part(rId)
                image_part.blob = new_pptx_img._blob
            except:
                pass
            k += 1


# def chaneg_text_of_slide(prs, new_texts, slide):
#     k = 0
    
#     for shape in prs.slides[int(slide)].shapes:
#         if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#             try:
#                 if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
#                     print(2)
#                     shape.text = "lol"
#                     print(f"TEXTBOX FOUND: text:{repr(shape.text/Pt(1))},width:{shape.width/Pt(1)}, heght:{shape.height/Pt(1)}, left:{shape.left/Pt(1)}, top:{shape.top/Pt(1)}")
#             except:
#                 pass
#                 # print("error")
#             k += 1

def replace_text(replacements: dict, slide):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    # slides = [slide for slide in prs.slides]
    shapes = []
    for shape in slide.shapes:
        shapes.append(shape)

    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if match in cell.text:
                            new_text = cell.text.replace(match, replacement)
                            cell.text = new_text


def change_image_and_text_on_slide(prs, slide, new_images, new_texts):
    # print(slide)
    change_image_of_slide(new_images, slide)
    replace_text(new_texts, slide)
    return prs

# prs = Presentation(r"C:\Users\Bogdan\java\pyth\t2p\presentations\doxy.pptx")
# # change_image_of_slide(prs, [r"C:\Users\Bogdan\java\pyth\t2p\presentations\server\imgs\test1\frog1.png"], 1)
# # replace_text({'Hello World213!': 'Hello!'}, prs.slides[1]) 

# change_image_and_text_on_slide(prs, 0, [r"C:\Users\Bogdan\java\pyth\t2p\presentations\server\imgs\test1\frog1.png"], {'Hello World213!': 'dfs!'})

# prs.save('doxy.pptx')

